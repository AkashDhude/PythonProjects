
import collections
import collections.abc
"""
Collection module is imported because of following error:-----
Container = collections.Container
AttributeError: module 'collections' has no attribute 'Container'
"""
from pptx import Presentation
from pptx.util import Inches, Pt

def addPPT(name = None,company_name=None, position_of_employee=None, email=None, mobile_num=None, website=None, fontname=None, address= None, file_upload = None, file_download = None):

    ppt = Presentation() # object creation
    # Creation of slide in ppt
    ppt.slide_width = Inches(3.5)
    ppt.slide_height = Inches(2)
    blank_slide_layout = ppt.slide_layouts[6] # creating a blank layout in ppt
    slide = ppt.slides.add_slide(blank_slide_layout) # creation of slide

    # Creation of textbox
    textbox_name = slide.shapes.add_textbox(left=Inches(0.125), top=Inches(0.15), width=Inches(2.25), height=Inches(0.25))
    tf1 = textbox_name.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = name
    p1.font.size = Pt(8)
    p1.font.name = fontname

    p6 = tf1.add_paragraph()
    p6.text = ""
    p6.font.size = Pt(2)

    p2 = tf1.add_paragraph()
    p2.text = position_of_employee
    p2.font.size = Pt(6)
    p2.font.name = fontname


    textbox_company_name = slide.shapes.add_textbox(left=Inches(0.125), top=Inches(0.9), width=Inches(2.5), height=Inches(0.25))
    tf3 = textbox_company_name.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = company_name
    p3.font.size = Pt(8)
    p3.font.bold = True
    p3.font.name = fontname

    textbox_address = slide.shapes.add_textbox(left=Inches(0.125), top=Inches(1.15), width=Inches(2.75), height=Inches(0.5))
    tf4 = textbox_address.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = address
    p4.font.size = Pt(6)

    p5 = tf4.add_paragraph()
    p5.text = f"Cell: +91 {mobile_num}"
    p5.font.size = Pt(6)
    p5.font.name = fontname

    textbox_email = slide.shapes.add_textbox(left=Inches(0.125), top=Inches(1.65), width=Inches(2.75), height=Inches(0.25))
    tf5 = textbox_email.text_frame
    p8 = tf5.paragraphs[0]
    p8.text = email + "  |  " + website
    p8.font.size = Pt(6)
    p8.font.name = fontname

    # Image setup
    image_path = file_upload
    slide.shapes.add_picture(image_path, left=Inches(2.5), top=Inches(0.2), height=Inches(0.7), width=Inches(0.8))

    # File save Location
    ppt.save(file_download + ".pptx")
